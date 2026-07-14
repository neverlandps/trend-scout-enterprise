"""PPTX report generation using python-pptx."""

import os
from datetime import datetime, timezone
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session

from trend_scout_enterprise.core.config import settings
from trend_scout_enterprise.models.models import Report
from trend_scout_enterprise.services.report_common import load_report_items, render_report_context


def _ensure_output_dir() -> str:
    os.makedirs(settings.output_dir, exist_ok=True)
    return settings.output_dir


def generate_ppt_report(db: Session, report: Report) -> str:
    """Generate a PPTX deck for the given report."""
    output_dir = _ensure_output_dir()
    item_ids = report.metadata_json.get("item_ids", [])
    filters = report.metadata_json.get("filters", {})
    items = load_report_items(db, report.owner_id, item_ids, filters)
    ctx = render_report_context(report.title, report.summary_text, items, filters)

    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = ctx["title"]
    slide.placeholders[1].text = f"Generated at {ctx['generated_at'][:19]}\nTotal items: {ctx['total_items']}"

    # Summary slide
    if ctx["summary"]:
        blank_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank_layout)
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_shape.text_frame.text = "Executive Summary"
        title_shape.text_frame.paragraphs[0].font.size = Pt(28)
        body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        body_shape.text_frame.text = ctx["summary"]

    # Item slides
    for item in items:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = item["title"]
        title_frame.paragraphs[0].font.size = Pt(24)

        details = (
            f"Source: {item['source_name']} ({item['source_type']})\n"
            f"Overall Score: {item['overall_score'] or 'N/A'}\n"
            f"Signal Strength: {item['signal_strength'] or 'N/A'}\n"
            f"URL: {item['url']}\n"
            f"Published: {item['published_at'][:10] if item['published_at'] else 'N/A'}"
        )
        detail_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(2.5))
        detail_box.text_frame.text = details
        detail_box.text_frame.paragraphs[0].font.size = Pt(14)

        summary_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.8))
        summary_box.text_frame.text = item["summary"]

    file_path = os.path.join(output_dir, f"{report.id}.pptx")
    prs.save(file_path)
    return file_path
